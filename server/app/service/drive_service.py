import os
import re
import time
import tempfile
import logging
from typing import Optional, Dict, List, Tuple
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload
from google.oauth2.credentials import Credentials
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from openpyxl import load_workbook


def _log(msg: str) -> None:
    """Force-flushed stdout print so messages survive worker SIGKILL."""
    print(f"[drive] {msg}", flush=True)

logger = logging.getLogger(__name__)


def _build_drive_service(credentials_dict: Dict):
    """Build Google Drive service from session credentials dict."""
    credentials = Credentials(
        token=credentials_dict.get('token'),
        refresh_token=credentials_dict.get('refresh_token'),
        token_uri=credentials_dict.get('token_uri'),
        client_id=credentials_dict.get('client_id'),
        client_secret=credentials_dict.get('client_secret'),
        scopes=credentials_dict.get('scopes'),
    )
    return build('drive', 'v3', credentials=credentials)


def _extract_folder_id(folder_url_or_id: str) -> str:
    """Extract folder ID from a Google Drive URL or return as-is if already an ID."""
    # Format: https://drive.google.com/drive/folders/FOLDER_ID
    match = re.search(r'/folders/([a-zA-Z0-9-_]+)', folder_url_or_id)
    if match:
        return match.group(1)
    return folder_url_or_id.strip()


def extract_docx_text(file_path: str) -> str:
    """Extract all text from a DOCX file (paragraphs and tables)."""
    try:
        doc = Document(file_path)
        text_parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)

        return '\n'.join(text_parts)
    except Exception as e:
        logger.error(f"Error extracting DOCX from {file_path}: {e}")
        return ''


def extract_markdown_text(file_path: str) -> str:
    """
    Read a text-based file (markdown, plain text, CSV, JSON, source code).
    Uses errors='replace' so files with non-UTF-8 bytes (common in real-world
    CSVs and old text files) still decode to *something* usable for retrieval.
    """
    try:
        with open(file_path, encoding='utf-8', errors='replace') as f:
            return f.read()
    except Exception as e:
        logger.error(f"Error reading text from {file_path}: {e}")
        return ''


def extract_pdf_text(file_path: str) -> str:
    """Extract text from a PDF file. Returns empty for image-only/scanned PDFs."""
    try:
        reader = PdfReader(file_path)
        pages = []
        for i, page in enumerate(reader.pages, start=1):
            page_text = page.extract_text() or ''
            if page_text.strip():
                pages.append(f"[Page {i}]\n{page_text}")
        return '\n\n'.join(pages)
    except Exception as e:
        logger.error(f"Error extracting PDF from {file_path}: {e}")
        return ''


def extract_xlsx_text(file_path: str) -> str:
    """Flatten every sheet's cells into 'col1 | col2 | col3' text rows for indexing."""
    try:
        wb = load_workbook(file_path, data_only=True, read_only=True)
        blocks: List[str] = []
        for sheet in wb.worksheets:
            blocks.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                line = ' | '.join('' if c is None else str(c) for c in row)
                if line.strip():
                    blocks.append(line)
        return '\n'.join(blocks)
    except Exception as e:
        logger.error(f"Error extracting XLSX from {file_path}: {e}")
        return ''


def extract_pptx_text(file_path: str) -> str:
    """
    Extract slide text + speaker notes from a .pptx file.
    Slides are tagged with [Slide N] and notes with [Notes] so the LLM can
    cite specific slides on retrieval.
    """
    try:
        prs = Presentation(file_path)
        slide_blocks: List[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            lines: List[str] = [f"[Slide {idx}]"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"[Notes] {notes}")
            slide_blocks.append('\n'.join(lines))
        return '\n\n'.join(slide_blocks)
    except Exception as e:
        logger.error(f"Error extracting PPTX from {file_path}: {e}")
        return ''


DOCX_MIME = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
PPTX_MIME = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
XLSX_MIME = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
PDF_MIME = 'application/pdf'

# Plain-text-style formats — all use the same passthrough extractor.
# Drive sometimes returns these MIME strings for source-code and config files.
_PLAIN_TEXT_MIMES = [
    'text/plain',
    'text/markdown',
    'text/csv',
    'text/tab-separated-values',
    'text/yaml',
    'application/yaml',
    'application/x-yaml',
    'application/json',
    'text/javascript',
    'application/javascript',
    'application/typescript',
    'text/x-python',
    'text/x-c',
    'text/x-c++',
    'text/x-java-source',
    'text/x-go',
    'text/x-rust',
    'application/xml',
    'text/xml',
    'application/sql',
    'text/x-sql',
]

# Extractor registry: MIME type -> extraction function.
# Keyed by the *effective* MIME type after any export.
EXTRACTORS = {
    DOCX_MIME: extract_docx_text,
    PPTX_MIME: extract_pptx_text,
    XLSX_MIME: extract_xlsx_text,
    PDF_MIME: extract_pdf_text,
    **{mime: extract_markdown_text for mime in _PLAIN_TEXT_MIMES},
}

# Google native MIME types are not downloadable directly; they must be exported
# via files().export_media(). Map each Google native type to the format we
# export it as, which must match a key in EXTRACTORS.
GOOGLE_EXPORTS = {
    'application/vnd.google-apps.document': 'text/markdown',
    'application/vnd.google-apps.presentation': PPTX_MIME,
    'application/vnd.google-apps.spreadsheet': 'text/csv',
}


def extract_text(file_path: str, mime_type: str) -> str:
    """Route to appropriate extractor based on MIME type."""
    extractor = EXTRACTORS.get(mime_type)
    if not extractor:
        logger.warning(f"No extractor for MIME type {mime_type}")
        return ''
    return extractor(file_path)


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    """Split text into overlapping chunks by character count."""
    if not text:
        return []

    chunks = []
    start = 0
    text_len = len(text)

    while start < text_len:
        end = min(start + chunk_size, text_len)
        chunk = text[start:end]

        if chunk.strip():
            chunks.append(chunk)

        start += chunk_size - overlap

    return chunks


class FolderAccessError(Exception):
    """Raised when a folder is not found, not accessible, or not a folder."""
    def __init__(self, code: str, message: str):
        super().__init__(message)
        self.code = code  # 'folder_not_found' | 'access_denied' | 'not_a_folder' | 'unknown'


def get_folder_name(folder_id: str, credentials_dict: Dict) -> str:
    """
    Fetch the human-readable name of a Drive folder.
    Raises FolderAccessError with a specific code on common failures.
    """
    try:
        service = _build_drive_service(credentials_dict)
        meta = service.files().get(fileId=folder_id, fields='name, mimeType').execute()
    except HttpError as e:
        status = e.resp.status if hasattr(e, 'resp') else None
        if status == 404:
            raise FolderAccessError('folder_not_found', "Folder not found. Check the URL.")
        if status in (401, 403):
            raise FolderAccessError('access_denied', "You don't have access to this folder.")
        raise FolderAccessError('unknown', f"Drive API error ({status}).")
    if meta.get('mimeType') != 'application/vnd.google-apps.folder':
        raise FolderAccessError('not_a_folder', "That URL points to a file, not a folder.")
    return meta.get('name') or 'Untitled folder'


def get_folder_files(folder_id: str, credentials_dict: Dict) -> List[Dict]:
    """List immediate children of a Drive folder (one level only)."""
    service = _build_drive_service(credentials_dict)
    files: List[Dict] = []
    page_token = None
    while True:
        results = service.files().list(
            q=f"'{folder_id}' in parents and trashed=false",
            spaces='drive',
            fields='files(id, name, mimeType, size, modifiedTime), nextPageToken',
            pageToken=page_token,
            pageSize=100,
        ).execute()
        files.extend(results.get('files', []))
        page_token = results.get('nextPageToken')
        if not page_token:
            break
    return files


def traverse_folder(folder_id: str, credentials_dict: Dict, max_depth: int = 5) -> Tuple[List[Dict], int]:
    """
    Recursively walk a folder up to `max_depth` levels deep. Returns
    (files, subfolder_count). Files are flat (no path tracking); subfolders are
    counted only for surfacing "we walked into N subfolders" in the UI.
    """
    FOLDER_MIME = 'application/vnd.google-apps.folder'
    all_files: List[Dict] = []
    subfolder_count = 0
    queue: List[Tuple[str, int]] = [(folder_id, 0)]
    visited: set = set()

    while queue:
        current_id, depth = queue.pop(0)
        if current_id in visited or depth > max_depth:
            continue
        visited.add(current_id)

        children = get_folder_files(current_id, credentials_dict)
        for c in children:
            if c.get('mimeType') == FOLDER_MIME:
                subfolder_count += 1
                queue.append((c['id'], depth + 1))
            else:
                all_files.append(c)

    return all_files, subfolder_count


def download_file(file_id: str, credentials_dict: Dict, mime_type: str) -> Optional[Tuple[str, str]]:
    """
    Download a file from Google Drive. For Google native types listed in
    GOOGLE_EXPORTS, export them as the target format instead. Returns
    (tmp_path, effective_mime_type) or None on failure.
    """
    try:
        service = _build_drive_service(credentials_dict)

        if mime_type in GOOGLE_EXPORTS:
            effective_mime = GOOGLE_EXPORTS[mime_type]
            request = service.files().export_media(fileId=file_id, mimeType=effective_mime)
        else:
            effective_mime = mime_type
            request = service.files().get_media(fileId=file_id)

        with tempfile.NamedTemporaryFile(delete=False, suffix='.tmp') as tmp_file:
            tmp_path = tmp_file.name

        with open(tmp_path, 'wb') as f:
            downloader = MediaIoBaseDownload(f, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()

        return tmp_path, effective_mime
    except Exception as e:
        logger.error(f"Error downloading file {file_id}: {e}")
        return None


def process_folder(folder_url_or_id: str, credentials_dict: Dict) -> Dict:
    """
    Process a Google Drive folder: list files (recursively), download supported
    files, extract text, chunk. Returns a structured result with success/error
    state, file count, chunk count, skipped files (with reasons), and
    subfolder count. The caller is responsible for storing chunks downstream.
    """
    try:
        t0 = time.time()
        folder_id = _extract_folder_id(folder_url_or_id)
        _log(f"START process_folder folder_id={folder_id}")

        try:
            folder_name = get_folder_name(folder_id, credentials_dict)
        except FolderAccessError as e:
            _log(f"folder access error: {e.code} — {e}")
            return {'status': 'error', 'error_code': e.code, 'message': str(e)}
        _log(f"folder_name={folder_name!r}")

        # Recursively gather files from this folder + all subfolders (up to max depth)
        t_list = time.time()
        files, subfolder_count = traverse_folder(folder_id, credentials_dict)
        _log(f"traverse took {time.time() - t_list:.2f}s, {len(files)} files, {subfolder_count} subfolders")

        supported_files = [
            f for f in files
            if f.get('mimeType') in EXTRACTORS or f.get('mimeType') in GOOGLE_EXPORTS
        ]
        unsupported_count = len(files) - len(supported_files)
        _log(f"{len(supported_files)} supported, {unsupported_count} unsupported")

        all_chunks: List[Dict] = []
        skipped_files: List[Dict] = []
        processed_count = 0

        for file_info in supported_files:
            file_id = file_info['id']
            file_name = file_info['name']
            mime_type = file_info['mimeType']
            _log(f"FILE start: {file_name} ({mime_type})")

            t_dl = time.time()
            download_result = download_file(file_id, credentials_dict, mime_type)
            _log(f"  download took {time.time() - t_dl:.2f}s")
            if not download_result:
                _log(f"  SKIP: download failed")
                skipped_files.append({'name': file_name, 'reason': 'download_failed'})
                continue
            tmp_path, effective_mime = download_result

            try:
                t_ex = time.time()
                text = extract_text(tmp_path, effective_mime)
                _log(f"  extract took {time.time() - t_ex:.2f}s, {len(text)} chars")
                if not text:
                    _log(f"  SKIP: no text")
                    skipped_files.append({'name': file_name, 'reason': 'no_text'})
                    continue

                text_chunks = chunk_text(text)
                if not text_chunks:
                    skipped_files.append({'name': file_name, 'reason': 'no_chunks'})
                    continue

                for chunk_index, chunk_text_content in enumerate(text_chunks):
                    all_chunks.append({
                        'document_text': chunk_text_content,
                        'metadata': {
                            'file_name': file_name,
                            'file_id': file_id,
                            'mime_type': mime_type,
                            'chunk_index': chunk_index,
                            'folder_id': folder_id,
                        }
                    })

                processed_count += 1
                _log(f"  FILE done: {len(text_chunks)} chunks")

            finally:
                try:
                    os.unlink(tmp_path)
                except Exception as e:
                    logger.warning(f"Failed to clean up temp file {tmp_path}: {e}")

        _log(f"END process_folder: {processed_count} files, {len(all_chunks)} chunks, "
             f"total {time.time() - t0:.2f}s")

        return {
            'status': 'success',
            'folder_id': folder_id,
            'folder_name': folder_name,
            'file_count': processed_count,
            'chunk_count': len(all_chunks),
            'chunks': all_chunks,
            'skipped_files': skipped_files,
            'unsupported_file_count': unsupported_count,
            'subfolder_count': subfolder_count,
            'files': [
                {
                    'name': f.get('name', ''),
                    'mime_type': f.get('mimeType', ''),
                    'size': int(f.get('size', 0) or 0),
                    'modified_time': f.get('modifiedTime', ''),
                }
                for f in files
            ],
        }

    except Exception as e:
        logger.error(f"Error processing folder {folder_url_or_id}: {e}")
        return {'status': 'error', 'error_code': 'unknown', 'message': str(e)}
