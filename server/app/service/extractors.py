"""
Text extraction from supported file formats + the EXTRACTORS / GOOGLE_EXPORTS
registries. Adding a new format means writing one extract_<x>_text function
and adding it to EXTRACTORS — no other code in the app changes.
"""
import logging
from typing import List

from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from openpyxl import load_workbook

logger = logging.getLogger(__name__)


def extract_docx_text(file_path: str) -> str:
    """Extract paragraph text + table cell text from a .docx file."""
    try:
        doc = Document(file_path)
        text_parts: List[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)
        return '\n'.join(text_parts)
    except Exception as e:
        logger.error(f"Error extracting DOCX from {file_path}: {e}")
        return ''


def extract_markdown_text(file_path: str) -> str:
    """
    Read a text-based file (markdown, plain text, CSV, JSON, source code).
    Uses errors='replace' so files with non-UTF-8 bytes (common in real-world
    CSVs and old text files) still decode to *something* usable for retrieval.
    """
    try:
        with open(file_path, encoding='utf-8', errors='replace') as f:
            return f.read()
    except Exception as e:
        logger.error(f"Error reading text from {file_path}: {e}")
        return ''


def extract_pdf_text(file_path: str) -> str:
    """Extract text from a PDF file. Returns empty for image-only/scanned PDFs."""
    try:
        reader = PdfReader(file_path)
        pages = []
        for i, page in enumerate(reader.pages, start=1):
            page_text = page.extract_text() or ''
            if page_text.strip():
                pages.append(f"[Page {i}]\n{page_text}")
        return '\n\n'.join(pages)
    except Exception as e:
        logger.error(f"Error extracting PDF from {file_path}: {e}")
        return ''


def extract_pptx_text(file_path: str) -> str:
    """
    Extract slide text + speaker notes from a .pptx file.
    Slides are tagged with [Slide N] and notes with [Notes] so the LLM can
    cite specific slides on retrieval.
    """
    try:
        prs = Presentation(file_path)
        slide_blocks: List[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            lines: List[str] = [f"[Slide {idx}]"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"[Notes] {notes}")
            slide_blocks.append('\n'.join(lines))
        return '\n\n'.join(slide_blocks)
    except Exception as e:
        logger.error(f"Error extracting PPTX from {file_path}: {e}")
        return ''


def extract_xlsx_text(file_path: str) -> str:
    """Flatten every sheet's cells into 'col1 | col2 | col3' text rows for indexing."""
    try:
        wb = load_workbook(file_path, data_only=True, read_only=True)
        blocks: List[str] = []
        for sheet in wb.worksheets:
            blocks.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                line = ' | '.join('' if c is None else str(c) for c in row)
                if line.strip():
                    blocks.append(line)
        return '\n'.join(blocks)
    except Exception as e:
        logger.error(f"Error extracting XLSX from {file_path}: {e}")
        return ''


# MIME type constants — used by both registries and by routes for filtering
DOCX_MIME = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
PPTX_MIME = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
XLSX_MIME = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
PDF_MIME = 'application/pdf'

# Plain-text-style formats — all use the same passthrough extractor.
# Drive sometimes returns these MIME strings for source-code and config files.
_PLAIN_TEXT_MIMES = [
    'text/plain',
    'text/markdown',
    'text/csv',
    'text/tab-separated-values',
    'text/yaml',
    'application/yaml',
    'application/x-yaml',
    'application/json',
    'text/javascript',
    'application/javascript',
    'application/typescript',
    'text/x-python',
    'text/x-c',
    'text/x-c++',
    'text/x-java-source',
    'text/x-go',
    'text/x-rust',
    'application/xml',
    'text/xml',
    'application/sql',
    'text/x-sql',
]

# Extractor registry: MIME type -> extraction function.
# Keyed by the *effective* MIME type after any export.
EXTRACTORS = {
    DOCX_MIME: extract_docx_text,
    PPTX_MIME: extract_pptx_text,
    XLSX_MIME: extract_xlsx_text,
    PDF_MIME: extract_pdf_text,
    **{mime: extract_markdown_text for mime in _PLAIN_TEXT_MIMES},
}

# Google native MIME types are not downloadable directly; they must be exported
# via files().export_media(). Map each Google native type to the format we
# export it as, which must match a key in EXTRACTORS.
GOOGLE_EXPORTS = {
    'application/vnd.google-apps.document': 'text/markdown',
    'application/vnd.google-apps.presentation': PPTX_MIME,
    'application/vnd.google-apps.spreadsheet': 'text/csv',
}


def extract_text(file_path: str, mime_type: str) -> str:
    """Route to the appropriate extractor based on MIME type."""
    extractor = EXTRACTORS.get(mime_type)
    if not extractor:
        logger.warning(f"No extractor for MIME type {mime_type}")
        return ''
    return extractor(file_path)
